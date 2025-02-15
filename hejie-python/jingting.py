from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import os

def personalize_slide(slide, name):
    keyword = "同学"
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # 检查文本中是否包含关键字
                    if keyword in run.text:
                        # 只替换文本，不改变其他格式
                        run.text = run.text.replace(keyword, f"{name}同学")
                        # 保留原有格式
                        run.font.size = Pt(16)  # 如果需要调整大小
                        paragraph.alignment = PP_ALIGN.LEFT  # 如果需要调整对齐

# 文件路径
pptx_template_path = 'data/tmp.pptx'
names_file_path = 'data/name.txt'
output_folder = 'output'

# 创建输出文件夹
if not os.path.exists(output_folder):
    os.makedirs(output_folder)

# 读取名字列表
with open(names_file_path, 'r', encoding='utf-8') as f:
    names = f.read().splitlines()

# 遍历名字列表
for name in names:
    new_prs = Presentation(pptx_template_path)
    slide = new_prs.slides[0]
    personalize_slide(slide, name)
    pptx_output_path = f'{output_folder}/{name}.pptx'
    new_prs.save(pptx_output_path)